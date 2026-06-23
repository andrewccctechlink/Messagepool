"""
Quattro Message Pool — File Parser (main dispatcher)
"""
import os
from typing import Optional


def parse_pdf(filepath: str) -> str:
    import pdfplumber
    parts = []
    with pdfplumber.open(filepath) as pdf:
        for i, page in enumerate(pdf.pages):
            t = page.extract_text()
            if t:
                parts.append(f"--- Page {i+1} ---\n{t}")
            for table in page.extract_tables():
                if table:
                    for row in table:
                        parts.append(" | ".join(str(c or "") for c in row))
    if not parts:
        raise ValueError("No text in PDF (may be image-based).")
    return "\n\n".join(parts)


def parse_excel(filepath: str) -> str:
    ext = os.path.splitext(filepath)[1].lower()
    if ext == ".xls":
        return _parse_xls(filepath)
    else:
        return _parse_xlsx(filepath)


def _parse_xlsx(filepath: str) -> str:
    import openpyxl
    parts = []
    wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
    for name in wb.sheetnames:
        ws = wb[name]
        parts.append(f"=== Sheet: {name} ===")
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() if c is not None else "" for c in row]
            if any(cells):
                rows.append(cells)
        if rows:
            parts.append("Columns: " + " | ".join(rows[0]))
            for r in rows[1:]:
                pairs = [f"{rows[0][i]}: {r[i]}" for i in range(min(len(rows[0]), len(r))) if r[i]]
                if pairs:
                    parts.append("; ".join(pairs))
    wb.close()
    return "\n".join(parts)


def _parse_xls(filepath: str) -> str:
    import xlrd
    parts = []
    wb = xlrd.open_workbook(filepath)
    for sheet in wb.sheets():
        parts.append(f"=== Sheet: {sheet.name} ===")
        rows = []
        for rx in range(sheet.nrows):
            cells = [str(sheet.cell_value(rx, cx)).strip() for cx in range(sheet.ncols)]
            if any(cells):
                rows.append(cells)
        if rows:
            parts.append("Columns: " + " | ".join(rows[0]))
            for r in rows[1:]:
                pairs = [f"{rows[0][i]}: {r[i]}" for i in range(min(len(rows[0]), len(r))) if r[i]]
                if pairs:
                    parts.append("; ".join(pairs))
    return "\n".join(parts)


def parse_csv(filepath: str) -> str:
    import csv
    enc = "utf-8"
    try:
        open(filepath, "r", encoding="utf-8").read(1024)
    except UnicodeDecodeError:
        enc = "gbk"
    with open(filepath, "r", encoding=enc, errors="replace") as f:
        reader = csv.reader(f)
        rows = list(reader)
    if not rows:
        raise ValueError("CSV is empty.")
    header = rows[0]
    parts = ["Columns: " + " | ".join(header)]
    for r in rows[1:]:
        pairs = [f"{header[i]}: {r[i].strip()}" for i in range(min(len(header), len(r))) if r[i].strip()]
        if pairs:
            parts.append("; ".join(pairs))
    return "\n".join(parts)


def parse_pptx(filepath: str) -> str:
    from pptx import Presentation
    prs = Presentation(filepath)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        if texts:
            parts.append(f"--- Slide {i} ---\n" + "\n".join(texts))
    return "\n\n".join(parts) if parts else "No text found in PPTX."


def parse_docx(filepath: str) -> str:
    from docx import Document
    doc = Document(filepath)
    parts = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            parts.append(t)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts) if parts else "No text found in DOCX."


def parse_file(filepath: str, api_key: str = None) -> tuple:
    """
    Parse any supported file. Returns (text, source_type).
    For images, api_key is needed for Gemini Vision.
    For .eml/.msg, also parses all attachments.
    """
    if not os.path.exists(filepath):
        raise FileNotFoundError(f"File not found: {filepath}")

    ext = os.path.splitext(filepath)[1].lower()

    if ext == ".pdf":
        return parse_pdf(filepath), "document"
    elif ext in (".xlsx", ".xls"):
        return parse_excel(filepath), "document"
    elif ext == ".csv":
        return parse_csv(filepath), "document"
    elif ext == ".pptx":
        return parse_pptx(filepath), "document"
    elif ext == ".docx":
        return parse_docx(filepath), "document"
    elif ext in (".jpg", ".jpeg", ".png"):
        from image_parser import parse_image
        return parse_image(filepath, api_key), "document"
    elif ext == ".eml":
        from email_parser import parse_eml_full
        return parse_eml_full(filepath, api_key), "email"
    elif ext == ".msg":
        from email_parser import parse_msg_full
        return parse_msg_full(filepath, api_key), "email"
    elif ext in (".txt", ".text", ".log"):
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            return f.read(), "document"
    else:
        raise ValueError(f"Unsupported: {ext}. Supported: pdf/xlsx/csv/pptx/docx/jpg/png/eml/msg/txt")
